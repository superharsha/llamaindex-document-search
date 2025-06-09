import streamlit as st
import asyncio
import os
import time
import glob
import tempfile
from io import BytesIO

from llama_index.core import (
    SimpleDirectoryReader,
    Settings,
    PromptTemplate,
    get_response_synthesizer,
)
from llama_index.core.node_parser import HierarchicalNodeParser, get_leaf_nodes
from llama_index.core.storage.docstore import SimpleDocumentStore
from llama_index.retrievers.bm25 import BM25Retriever
from llama_index.llms.google_genai import GoogleGenAI
from llama_index.core.query_engine import RetrieverQueryEngine, TransformQueryEngine
from llama_index.core.indices.query.query_transform import HyDEQueryTransform
from dotenv import load_dotenv

# PDF conversion imports
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet

# Load environment variables for local development
load_dotenv()

# Function to get secrets from either Streamlit secrets or environment variables
def get_secret(key):
    try:
        # Try Streamlit secrets first (for deployment)
        return st.secrets[key]
    except (KeyError, AttributeError):
        # Fall back to environment variables (for local development)
        return os.getenv(key)

# Set Gemini API key
gemini_api_key = get_secret("GEMINI_API_KEY")
if gemini_api_key:
    os.environ["GOOGLE_API_KEY"] = gemini_api_key

# ────────────────────────────────────────────────────────────────────────────────
# 1. Build (and cache) a local, BM25-only query engine with HyDE + Gemini
# ────────────────────────────────────────────────────────────────────────────────
# Dependencies: python-pptx, Pillow, openpyxl, docx2txt (lightweight, no torch/transformers)

@st.cache_resource
def build_local_query_engine(sim_top_k: int = 5):
    """
    1) Read all files under ./my_local_docs/
    2) Chunk them hierarchically (simplified for performance)
    3) Populate an in-memory SimpleDocumentStore
    4) Build BM25Retriever → HyDE → Gemini
    """
    # Show progress
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    # 1) HyDE transformer (to rewrite queries + include original)
    hyde = HyDEQueryTransform(include_original=True)
    status_text.text("Initializing query transformer...")
    progress_bar.progress(10)

    # 2) Read & chunk documents with error handling
    status_text.text("Loading documents...")
    documents = []
    doc_dir = "./my_local_docs/"
    failed_files = []
    successful_files = []
    
    # Get list of files first
    file_paths = [f for f in glob.glob(doc_dir + "*") if os.path.isfile(f)]
    total_files = len(file_paths)
    
    def convert_docx_to_pdf(docx_path):
        """Convert DOCX file to PDF and return the PDF content as bytes"""
        try:
            doc = Document(docx_path)
            
            # Create PDF in memory
            pdf_buffer = BytesIO()
            pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Extract text from DOCX and add to PDF
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    p = Paragraph(paragraph.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            pdf_doc.build(story)
            pdf_buffer.seek(0)
            return pdf_buffer.getvalue()
        except Exception as e:
            st.warning(f"Failed to convert DOCX to PDF: {e}")
            return None

    def convert_pptx_to_pdf(pptx_path):
        """Convert PPTX file to PDF and return the PDF content as bytes"""
        try:
            prs = Presentation(pptx_path)
            
            # Create PDF in memory
            pdf_buffer = BytesIO()
            pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Extract text from PPTX slides and add to PDF
            for i, slide in enumerate(prs.slides):
                story.append(Paragraph(f"Slide {i+1}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        p = Paragraph(shape.text, styles['Normal'])
                        story.append(p)
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 24))
            
            pdf_doc.build(story)
            pdf_buffer.seek(0)
            return pdf_buffer.getvalue()
        except Exception as e:
            st.warning(f"Failed to convert PPTX to PDF: {e}")
            return None

    # Try loading files individually with conversion for DOCX/PPTX
    for i, file_path in enumerate(file_paths):
        try:
            # Update progress
            progress_bar.progress(10 + (i * 30 // total_files))
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_name)[1].lower()
            status_text.text(f"Processing {file_name}... ({i+1}/{total_files})")
            
            if file_ext in ['.docx']:
                # Convert DOCX to PDF
                status_text.text(f"Converting DOCX to PDF: {file_name}...")
                pdf_content = convert_docx_to_pdf(file_path)
                if pdf_content:
                    # Create temporary PDF file
                    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                        temp_pdf.write(pdf_content)
                        temp_pdf_path = temp_pdf.name
                    
                    # Load the converted PDF
                    file_docs = SimpleDirectoryReader(input_files=[temp_pdf_path]).load_data()
                    # Clean up temp file
                    os.unlink(temp_pdf_path)
                    
                    documents.extend(file_docs)
                    successful_files.append(f"{file_name} (converted to PDF)")
                else:
                    failed_files.append(f"{file_name}: DOCX conversion failed")
                    
            elif file_ext in ['.pptx']:
                # Convert PPTX to PDF
                status_text.text(f"Converting PPTX to PDF: {file_name}...")
                pdf_content = convert_pptx_to_pdf(file_path)
                if pdf_content:
                    # Create temporary PDF file
                    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                        temp_pdf.write(pdf_content)
                        temp_pdf_path = temp_pdf.name
                    
                    # Load the converted PDF
                    file_docs = SimpleDirectoryReader(input_files=[temp_pdf_path]).load_data()
                    # Clean up temp file
                    os.unlink(temp_pdf_path)
                    
                    documents.extend(file_docs)
                    successful_files.append(f"{file_name} (converted to PDF)")
                else:
                    failed_files.append(f"{file_name}: PPTX conversion failed")
                    
            else:
                # Process other files normally (PDF, XLSX, images, etc.)
                file_docs = SimpleDirectoryReader(input_files=[file_path]).load_data()
                documents.extend(file_docs)
                successful_files.append(file_name)
                
        except Exception as file_error:
            failed_files.append(f"{file_name}: {str(file_error)}")
            continue
    
    progress_bar.progress(40)
    status_text.text("Documents loaded. Processing chunks...")
    
    # Show loading results
    if successful_files:
        st.success(f"✅ Successfully loaded {len(successful_files)} files")
    if failed_files:
        st.warning(f"❌ Failed to load {len(failed_files)} files (likely corrupted)")
        with st.expander("View failed files"):
            for failed_file in failed_files:
                st.text(failed_file)
    
    if not documents:
        raise RuntimeError("No files could be loaded from ./my_local_docs/. Please check your files.")

    # Simplified hierarchical parsing for better performance
    status_text.text("Creating document chunks...")
    parser = HierarchicalNodeParser.from_defaults(
        chunk_sizes=[4096, 2048, 1024, 512]  # Reverted to original chunk sizes
    )
    nodes = parser.get_nodes_from_documents(documents)
    leaf_nodes = get_leaf_nodes(nodes)
    
    progress_bar.progress(60)
    status_text.text("Building search index...")

    # 3) Build in-memory BM25 index
    docstore = SimpleDocumentStore()
    docstore.add_documents(nodes)

    bm25_retriever = BM25Retriever.from_defaults(
        docstore=docstore,
        similarity_top_k=sim_top_k,
    )
    
    progress_bar.progress(80)
    status_text.text("Initializing AI model...")

    # 4) Attach Gemini as the final answer generator
    llm = GoogleGenAI(
        model="gemini-2.0-flash",
        temperature=0,
        api_key=gemini_api_key
    )
    Settings.llm = llm

    # 5) Build a RetrieverQueryEngine (BM25 → simple synthesizer)
    primitive_engine = RetrieverQueryEngine.from_args(
        bm25_retriever,
        response_synthesizer=get_response_synthesizer(
            text_qa_template=PromptTemplate(
                """You are an expert document analyst. Use the provided context to answer the user's question accurately and comprehensively.

CONTEXT INFORMATION:
{context_str}

INSTRUCTIONS:
- Answer based ONLY on the information provided in the context above
- If the context doesn't contain enough information to fully answer the question, clearly state what information is missing
- Provide specific details, quotes, and examples from the context when relevant
- Structure your response clearly with proper formatting
- If multiple documents are referenced, distinguish between them in your answer
- If you find contradictory information, acknowledge and explain the discrepancies

QUESTION: {query_str}

ANSWER:"""
            ),
            response_mode="compact",
        ),
    )

    # 6) Wrap with HyDE → Gemini
    engine = TransformQueryEngine(primitive_engine, hyde)
    
    progress_bar.progress(100)
    status_text.text("✅ Index ready!")
    
    # Clear progress indicators after a moment
    time.sleep(1)
    progress_bar.empty()
    status_text.empty()
    
    return engine

# Helper to call the query engine synchronously
def run_query(engine, query: str):
    return engine.query(query)


# ────────────────────────────────────────────────────────────────────────────────
# 2. Streamlit UI
# ────────────────────────────────────────────────────────────────────────────────

st.set_page_config(page_title="Local BM25 + HyDE Query", layout="wide")
st.title("📚 Local Document Search (BM25 + HyDE → Gemini)")

st.markdown("""
This app:
1. Reads **`./my_local_docs/`** (supports `.txt`, `.pdf`, `.docx`)  
2. Chunks everything into a BM25 index  
3. Wraps BM25 with HyDE → Gemini for answer generation  
4. Let's you type a query and see:
   - **Answer** (generated by Gemini)  
   - **Top matching chunks** (with BM25 scores)
""")

# Performance warning
doc_count = len([f for f in glob.glob("./my_local_docs/*") if os.path.isfile(f)])
if doc_count > 30:
    st.warning(f"⚠️ **Performance Notice**: You have {doc_count} documents. Initial indexing may take 1-2 minutes on Streamlit Cloud due to resource limitations. The index is cached after first run.")
else:
    st.info(f"📁 Found {doc_count} documents ready for search.")

query = st.text_input("Enter your query here:", placeholder="e.g. What are the main points in document X?")

if st.button("🔍 Run Query") and query.strip():
    try:
        with st.spinner("🔄 Building index (if first run) and querying…"):
            # Start timing
            start_time = time.time()
            
            engine = build_local_query_engine(sim_top_k=20)
            resp = run_query(engine, query)
            
            # End timing
            end_time = time.time()
            elapsed_time = end_time - start_time
            
    except Exception as e:
        st.error(f"Error: {e}")
    else:
        # Display timing information
        st.success(f"⏱️ **Query completed in {elapsed_time:.2f} seconds**")
        
        st.subheader("🗒️ Answer")
        st.write(resp.response)

        st.subheader("🔝 Top Chunks")
        for node_info in resp.source_nodes:
            txt = node_info.node.text.replace("\n", " ").strip()
            score = round(node_info.score, 3)
            # Show the first 200 characters of each chunk
            display_text = txt[:200] + ("…" if len(txt) > 200 else "")
            st.write(f"- {display_text}  _(score: {score})_")
